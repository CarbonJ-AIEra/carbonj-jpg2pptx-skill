#!/usr/bin/env python3
import json
import sys
import zipfile
from pathlib import Path

from pptx import Presentation


def main() -> int:
    if len(sys.argv) != 2:
        print("Usage: check_editability.py <deck.pptx>", file=sys.stderr)
        return 2
    deck = Path(sys.argv[1]).expanduser().resolve()
    if not deck.is_file():
        print(f"File not found: {deck}", file=sys.stderr)
        return 2
    if not zipfile.is_zipfile(deck):
        print(f"Not a valid ZIP-based PPTX: {deck}", file=sys.stderr)
        return 1

    prs = Presentation(str(deck))
    slides = []
    total_text = 0
    total_images = 0
    total_shapes = 0
    for slide_no, slide in enumerate(prs.slides, 1):
        text_objects = []
        images = 0
        for shape in slide.shapes:
            total_shapes += 1
            if shape.shape_type == 13:
                images += 1
                total_images += 1
            if getattr(shape, "has_text_frame", False):
                text = " ".join(shape.text.split())
                if text:
                    text_objects.append(text)
                    total_text += 1
        slides.append({
            "slide": slide_no,
            "shape_count": len(slide.shapes),
            "text_object_count": len(text_objects),
            "image_count": images,
            "texts": text_objects,
        })

    result = {
        "path": str(deck),
        "slide_count": len(prs.slides),
        "shape_count": total_shapes,
        "text_object_count": total_text,
        "image_count": total_images,
        "slides": slides,
    }
    print(json.dumps(result, ensure_ascii=False, indent=2))
    if len(prs.slides) == 0 or total_shapes == 0:
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

